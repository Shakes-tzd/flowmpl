#!/usr/bin/env python3
"""Generate the tariff carve-out scene as a native PowerPoint file.

Reads ``pptx_preferences.json`` for user-corrected positions (written by
``sync_pptx_feedback.py`` after each editing session).  Falls back to
``DEFAULTS`` for any key not yet in preferences — so deleting the prefs
file safely reverts to the first-run layout.

After generation, writes ``pptx_baseline.json`` so subsequent edits can
be diffed against it by ``sync_pptx_feedback.py``.

Usage
-----
    uv run --with python-pptx notebooks/generate_tariff_scene.py
    # → notebooks/test_tariff_scene.pptx

Active-learning workflow
------------------------
    1.  Run this script  →  test_tariff_scene.pptx
    2.  Open in PowerPoint, adjust positions / sizes / styles
    3.  Save and close PowerPoint
    4.  uv run --with python-pptx notebooks/sync_pptx_feedback.py
    5.  Repeat from step 1 — positions improve each cycle

Design tokens (flowmpl design.py)
----------------------------------
    CONCEPT_WHITE  #FFFFFF    slide background (default layout)
    CONCEPT_INK    #1A1A1A    text, borders
    CONCEPT_YELLOW #F5C842    accent bars, moneybag/coins icons
    CONCEPT_MUTED  #888888    secondary labels
    figsize        (12, 6.75)  ↔  slide 13.333″ × 7.5″  (same 16:9)
"""

from __future__ import annotations

import io
import json
from pathlib import Path

import numpy as np
from PIL import Image as _PIL
from pptx import Presentation
from pptx.chart.data import ChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

# ── Paths ──────────────────────────────────────────────────────────────────────
_HERE       = Path(__file__).parent
ICONS_DIR   = _HERE / "assets" / "icons"
PREFS_FILE  = _HERE / "pptx_preferences.json"
BASELINE_FILE = _HERE / "pptx_baseline.json"
OUTPUT_FILE = _HERE / "test_tariff_scene.pptx"

# ── Design tokens ──────────────────────────────────────────────────────────────
INK      = RGBColor(0x1A, 0x1A, 0x1A)
YELLOW   = RGBColor(0xF5, 0xC8, 0x42)
WHITE    = RGBColor(0xFF, 0xFF, 0xFF)
MUTED    = RGBColor(0x88, 0x88, 0x88)
GRAY_BAR = RGBColor(0xD0, 0xD0, 0xD0)  # pre-carve-out bar
PANEL_BG = RGBColor(0xF7, 0xF3, 0xEE)  # warm cream
ARROW_C  = RGBColor(0xCC, 0xCC, 0xAA)  # muted warm gray arrow

# ── Slide dimensions: 16:9 widescreen = flowmpl figsize (12, 6.75) ────────────
SW, SH = 13.333, 7.5          # inches
SW_EMU = int(SW * 914400)
SH_EMU = int(SH * 914400)

# ── Auto-shape type integers ───────────────────────────────────────────────────
_RECT        = 1    # MSO_AUTO_SHAPE_TYPE.RECTANGLE
_RIGHT_ARROW = 13   # MSO_AUTO_SHAPE_TYPE.RIGHT_ARROW

# ── Defaults: the first-run layout (before any user corrections) ───────────────
DEFAULTS: dict = {
    "style": {
        "background":       "split",      # "split" or "full_warm"
        "background_color": "#F7F3EE",
        "arrow_style":      "connector",  # "connector" or "block"
    },
    "chart": {"x0": 0.12, "y0": 0.08, "x1": 0.52, "y1": 0.92},
    "arrow_connector": {
        "tail":     [0.18, 0.18],
        "tip":      [0.55, 0.80],
        "color":    "#CCCCAA",
        "width_pt": 22,
    },
    "arrow_block": {
        "x0": 0.326, "y0": 0.606, "x1": 0.611, "y1": 0.699,
        "rotation": 302.57,
    },
    "icons": {
        "server_left":       {"x0": 0.00, "y0": 0.73, "x1": 0.12, "y1": 0.93},
        "cpu":               {"x0": 0.00, "y0": 0.47, "x1": 0.12, "y1": 0.67},
        "database":          {"x0": 0.00, "y0": 0.08, "x1": 0.12, "y1": 0.28},
        "server_top_center": {"x0": 0.20, "y0": 0.79, "x1": 0.32, "y1": 0.97},
        "factory":           {"x0": 0.28, "y0": 0.02, "x1": 0.54, "y1": 0.22},
        "coins":             {"x0": 0.49, "y0": 0.40, "x1": 0.62, "y1": 0.58},
        "shield":            {"x0": 0.50, "y0": 0.22, "x1": 0.61, "y1": 0.38},
        "cloud":             {"x0": 0.70, "y0": 0.84, "x1": 0.82, "y1": 0.97},
        "bank":              {"x0": 0.82, "y0": 0.76, "x1": 0.96, "y1": 0.95},
        "moneybag_1":        {"x0": 0.33, "y0": 0.80, "x1": 0.50, "y1": 0.93},
        "moneybag_2":        {"x0": 0.45, "y0": 0.59, "x1": 0.60, "y1": 0.88},
    },
    "text": {
        "title":      {"x0": 0.57, "y0": 0.68, "x1": 0.97, "y1": 0.86},
        "subhead":    {"x0": 0.57, "y0": 0.56, "x1": 0.97, "y1": 0.68},
        "body":       {"x0": 0.57, "y0": 0.28, "x1": 0.97, "y1": 0.56},
        "subsidy":    {"x0": 0.57, "y0": 0.21, "x1": 0.97, "y1": 0.28},
        "axis_label": {"x0": 0.01, "y0": 0.88, "x1": 0.10, "y1": 0.99},
    },
}

# Icons: (logical_name, svg_stem, ink_color_hex)
# Moneybag and coins use CONCEPT_YELLOW to match the flowmpl notebook.
ICON_DEFS: list[tuple[str, str, str]] = [
    ("server_left",       "server",    "#1A1A1A"),
    ("cpu",               "cpu",       "#1A1A1A"),
    ("database",          "database",  "#1A1A1A"),
    ("server_top_center", "server",    "#1A1A1A"),
    ("factory",           "factory",   "#1A1A1A"),
    ("coins",             "coins",     "#F5C842"),
    ("shield",            "shield",    "#1A1A1A"),
    ("cloud",             "cloud",     "#1A1A1A"),
    ("bank",              "bank",      "#1A1A1A"),
    ("moneybag_1",        "moneybag",  "#F5C842"),
    ("moneybag_2",        "moneybag",  "#F5C842"),
]


# ── Preferences ────────────────────────────────────────────────────────────────

def _deep_merge(base: dict, override: dict) -> dict:
    """Recursively merge override into base; skip keys starting with '_'."""
    result = dict(base)
    for k, v in override.items():
        if k.startswith("_"):
            continue
        if isinstance(v, dict) and isinstance(result.get(k), dict):
            result[k] = _deep_merge(result[k], v)
        else:
            result[k] = v
    return result


def load_config() -> dict:
    """Merge DEFAULTS with saved preferences (preferences win on conflict)."""
    if PREFS_FILE.exists():
        with open(PREFS_FILE) as f:
            prefs = json.load(f)
        return _deep_merge(DEFAULTS, prefs)
    return dict(DEFAULTS)


# ── Coordinate helpers ─────────────────────────────────────────────────────────

def fr(x0: float, y0: float, x1: float, y1: float):
    """Axes-fraction bbox (matplotlib: y=0 bottom) → (left, top, w, h) EMU."""
    return (
        Inches(x0 * SW),
        Inches((1.0 - y1) * SH),   # y1 = matplotlib top → PPT top
        Inches((x1 - x0) * SW),
        Inches((y1 - y0) * SH),
    )


def pt_xy(x: float, y: float):
    """Single fraction point → (left, top) EMU."""
    return Inches(x * SW), Inches((1.0 - y) * SH)


# ── Shape baseline extraction (for sync diffing) ───────────────────────────────

def _shape_record(shape) -> dict:
    """Extract position/size/rotation from a shape as fraction coords."""
    x0 = round(shape.left / SW_EMU, 4)
    x1 = round((shape.left + shape.width) / SW_EMU, 4)
    y1 = round(1.0 - shape.top / SH_EMU, 4)            # ppt top → mpl top
    y0 = round(1.0 - (shape.top + shape.height) / SH_EMU, 4)
    rec = {"x0": x0, "y0": y0, "x1": x1, "y1": y1,
           "shape_type": str(shape.shape_type)}
    rot = getattr(shape, "rotation", None)
    if rot:
        rec["rotation"] = round(rot, 3)
    return rec


# ── Shape helpers ──────────────────────────────────────────────────────────────

def _solid(shape, color: RGBColor) -> None:
    shape.fill.solid()
    shape.fill.fore_color.rgb = color


def _no_line(shape) -> None:
    shape.line.color.rgb = WHITE
    shape.line.width = Pt(0.1)


def _add_text(
    slide,
    name: str,
    text: str,
    x0: float, y0: float, x1: float, y1: float,
    *,
    bold: bool = False,
    size: float = 10,
    color: RGBColor | None = None,
    align: PP_ALIGN = PP_ALIGN.LEFT,
):
    txb = slide.shapes.add_textbox(*fr(x0, y0, x1, y1))
    txb.name = name
    tf = txb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.bold = bold
    run.font.size = Pt(size)
    run.font.color.rgb = color or INK
    return txb


# ── Icon loading ───────────────────────────────────────────────────────────────

def _icon_png(stem: str, color: str = "#1A1A1A", size: int = 256) -> io.BytesIO:
    """Rasterise a local SVG icon and return PNG bytes (transparent bg)."""
    from flowmpl import load_icon   # noqa: PLC0415  (optional dep guard)
    arr = load_icon(ICONS_DIR / f"{stem}.svg", color=color, size=size)
    # arr is float32 RGBA in [0,1]; convert to uint8 PIL image → PNG bytes
    img = _PIL.fromarray((arr * 255).astype(np.uint8), "RGBA")
    buf = io.BytesIO()
    img.save(buf, format="PNG")
    buf.seek(0)
    return buf


def _add_icon(
    slide,
    name: str,
    stem: str,
    x0: float, y0: float, x1: float, y1: float,
    color: str = "#1A1A1A",
) -> None:
    """Embed a sketch SVG icon as a transparent PNG at fraction coordinates."""
    buf = _icon_png(stem, color=color)
    pic = slide.shapes.add_picture(buf, *fr(x0, y0, x1, y1))
    pic.name = name


# ── Main build ─────────────────────────────────────────────────────────────────

def build() -> None:
    cfg = load_config()
    prs = Presentation()
    prs.slide_width  = Inches(SW)
    prs.slide_height = Inches(SH)
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank

    style  = cfg["style"]
    bg_hex = style["background_color"].lstrip("#")
    bg_rgb = RGBColor(int(bg_hex[0:2], 16), int(bg_hex[2:4], 16), int(bg_hex[4:6], 16))

    # ── 1. Background ──────────────────────────────────────────────────────
    if style["background"] == "full_warm":
        # Single full-slide warm cream (learned preference from session 1)
        bg = slide.shapes.add_shape(_RECT, Inches(0), Inches(0), Inches(SW), Inches(SH))
        _solid(bg, bg_rgb)
        _no_line(bg)
        bg.name = "bg"
    else:
        # Original split: white base + warm right panel
        bg = slide.shapes.add_shape(_RECT, Inches(0), Inches(0), Inches(SW), Inches(SH))
        _solid(bg, WHITE)
        _no_line(bg)
        bg.name = "bg"
        rp = slide.shapes.add_shape(_RECT, *fr(0.54, 0.0, 1.0, 1.0))
        _solid(rp, bg_rgb)
        _no_line(rp)
        rp.name = "bg_right_panel"

    # ── 2. Arrow (drawn before icons so it sits behind them) ──────────────
    if style["arrow_style"] == "block":
        ab = cfg["arrow_block"]
        arrow = slide.shapes.add_shape(_RIGHT_ARROW, *fr(ab["x0"], ab["y0"], ab["x1"], ab["y1"]))
        arrow.rotation = ab["rotation"]
        _solid(arrow, ARROW_C)
        _no_line(arrow)
        arrow.name = "arrow"
    else:
        # Faded line connector
        ac = cfg["arrow_connector"]
        t_x, t_y = pt_xy(*ac["tail"])
        h_x, h_y = pt_xy(*ac["tip"])
        from pptx.enum.shapes import MSO_CONNECTOR_TYPE  # noqa: PLC0415
        conn = slide.shapes.add_connector(MSO_CONNECTOR_TYPE.STRAIGHT, t_x, t_y, h_x, h_y)
        hx = ac["color"].lstrip("#")
        conn.line.color.rgb = RGBColor(int(hx[0:2], 16), int(hx[2:4], 16), int(hx[4:6], 16))
        conn.line.width = Pt(ac["width_pt"])
        conn.name = "arrow"

    # ── 3. Bar chart — native, editable ───────────────────────────────────
    chart_data = ChartData()
    chart_data.categories = ["Pre-\nCarve-out", "Post-\nCarve-out"]
    chart_data.add_series("Relative Import Volume", (150, 450))
    c = cfg["chart"]
    chart_frame = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED, *fr(c["x0"], c["y0"], c["x1"], c["y1"]), chart_data
    )
    chart_frame.name = "chart"
    ch = chart_frame.chart
    ch.has_legend = False
    ch.plots[0].gap_width = 80
    ch.series[0].points[0].format.fill.solid()
    ch.series[0].points[0].format.fill.fore_color.rgb = GRAY_BAR
    ch.series[0].points[1].format.fill.solid()
    ch.series[0].points[1].format.fill.fore_color.rgb = YELLOW

    # ── 4. Icons (actual SVG sketch assets, transparent PNG) ──────────────
    for icon_name, stem, color in ICON_DEFS:
        b = cfg["icons"][icon_name]
        _add_icon(slide, f"icon_{icon_name}", stem, b["x0"], b["y0"], b["x1"], b["y1"], color=color)

    # ── 5. Text panel ──────────────────────────────────────────────────────
    t = cfg["text"]
    _add_text(slide, "text_title", "TARIFF CARVE-OUT",
              t["title"]["x0"], t["title"]["y0"], t["title"]["x1"], t["title"]["y1"],
              bold=True, size=20, color=INK)
    _add_text(slide, "text_subhead", "De-Risking Zone",
              t["subhead"]["x0"], t["subhead"]["y0"], t["subhead"]["x1"], t["subhead"]["y1"],
              bold=True, size=13, color=MUTED)
    _add_text(
        slide, "text_body",
        ("Solar and battery storage equipment\n"
         "exempted from retaliatory tariffs.\n\n"
         "Domestic manufacturing subsidies\n"
         "preserved through carve-out clause.\n\n"
         "Import volume projected to rise 3×\n"
         "within two policy cycles."),
        t["body"]["x0"], t["body"]["y0"], t["body"]["x1"], t["body"]["y1"],
        bold=False, size=11, color=INK)
    _add_text(slide, "text_subsidy", "Subsidy",
              t["subsidy"]["x0"], t["subsidy"]["y0"], t["subsidy"]["x1"], t["subsidy"]["y1"],
              bold=False, size=10, color=MUTED)
    _add_text(slide, "text_axis_label", "Relative\nImport\nVolume",
              t["axis_label"]["x0"], t["axis_label"]["y0"],
              t["axis_label"]["x1"], t["axis_label"]["y1"],
              bold=False, size=7, color=MUTED, align=PP_ALIGN.CENTER)

    # ── Save PPTX ──────────────────────────────────────────────────────────
    prs.save(str(OUTPUT_FILE))
    print(f"Saved → {OUTPUT_FILE}")

    # ── Save baseline for sync diffing ─────────────────────────────────────
    baseline = {}
    for shape in slide.shapes:
        if shape.name:
            baseline[shape.name] = _shape_record(shape)
    BASELINE_FILE.write_text(json.dumps(baseline, indent=2))
    print(f"Baseline → {BASELINE_FILE}")


if __name__ == "__main__":
    build()
