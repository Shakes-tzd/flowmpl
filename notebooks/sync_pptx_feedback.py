#!/usr/bin/env python3
"""Sync user edits from test_tariff_scene.pptx → pptx_preferences.json.

Reads the user-modified PPTX, diffs every named shape against
``pptx_baseline.json`` (what ``generate_tariff_scene.py`` last wrote),
then updates ``pptx_preferences.json`` with the corrected positions.

On the next generation run, the corrected positions become the starting
point — so the layout improves each cycle without the user re-making
the same corrections.

Usage
-----
Run this AFTER editing test_tariff_scene.pptx in PowerPoint:

    uv run --with python-pptx notebooks/sync_pptx_feedback.py

Insight extraction
------------------
Beyond raw coordinate capture, the script derives style preferences:

* Background style (split vs full_warm)
* Arrow style preference (connector vs block shape)
* Average icon scale factor (consistent sizing corrections)
* Per-session notes summarising what changed and why
"""

from __future__ import annotations

import json
import math
from datetime import date
from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

# ── Paths ──────────────────────────────────────────────────────────────────────
_HERE         = Path(__file__).parent
PPTX_FILE     = _HERE / "test_tariff_scene.pptx"
BASELINE_FILE = _HERE / "pptx_baseline.json"
PREFS_FILE    = _HERE / "pptx_preferences.json"

SW, SH = 13.333, 7.5
SW_EMU = int(SW * 914400)
SH_EMU = int(SH * 914400)


# ── Extraction ─────────────────────────────────────────────────────────────────

def _extract_shapes(pptx_path: Path) -> dict[str, dict]:
    """Return {shape_name: record} for every named shape in the first slide."""
    prs  = Presentation(str(pptx_path))
    slide = prs.slides[0]
    shapes = {}
    for shape in slide.shapes:
        if not shape.name:
            continue
        x0 = round(shape.left / SW_EMU, 4)
        x1 = round((shape.left + shape.width) / SW_EMU, 4)
        y1 = round(1.0 - shape.top / SH_EMU, 4)
        y0 = round(1.0 - (shape.top + shape.height) / SH_EMU, 4)
        rec: dict = {
            "x0": x0, "y0": y0, "x1": x1, "y1": y1,
            "shape_type": str(shape.shape_type),
        }
        rot = getattr(shape, "rotation", None)
        if rot:
            rec["rotation"] = round(rot, 3)
        # Fill color
        try:
            if shape.fill.type is not None:
                rec["fill"] = "#" + str(shape.fill.fore_color.rgb)
        except Exception:
            pass
        shapes[shape.name] = rec
    return shapes


# ── Diff ───────────────────────────────────────────────────────────────────────

_THRESH = 0.002   # fraction-coordinate change below this is considered noise

def _diff(baseline: dict, current: dict) -> dict[str, dict]:
    """Return changed/deleted/new shapes with old/new values and deltas."""
    changes: dict[str, dict] = {}

    for name, old in baseline.items():
        if name not in current:
            changes[name] = {"status": "deleted", "old": old}
            continue
        new = current[name]
        dx0 = round(new["x0"] - old["x0"], 4)
        dy0 = round(new["y0"] - old["y0"], 4)
        dx1 = round(new["x1"] - old["x1"], 4)
        dy1 = round(new["y1"] - old["y1"], 4)
        dw  = round((new["x1"] - new["x0"]) - (old["x1"] - old["x0"]), 4)
        dh  = round((new["y1"] - new["y0"]) - (old["y1"] - old["y0"]), 4)

        pos_changed  = any(abs(d) > _THRESH for d in (dx0, dy0, dx1, dy1))
        type_changed = new.get("shape_type") != old.get("shape_type")
        rot_old      = old.get("rotation", 0) or 0
        rot_new      = new.get("rotation", 0) or 0
        rot_changed  = abs(rot_new - rot_old) > 0.5
        fill_changed = new.get("fill") != old.get("fill") and new.get("fill") is not None

        if pos_changed or type_changed or rot_changed or fill_changed:
            changes[name] = {
                "status": "changed",
                "old": old,
                "new": new,
                "delta": {"dx0": dx0, "dy0": dy0, "dx1": dx1, "dy1": dy1,
                          "dw": dw, "dh": dh},
            }

    # Shapes added by user (in current but not baseline)
    for name in current:
        if name not in baseline:
            changes[name] = {"status": "added", "new": current[name]}

    return changes


# ── Insight derivation ─────────────────────────────────────────────────────────

def _derive_insights(changes: dict, current: dict) -> dict:
    """Extract style preferences from the pattern of corrections."""
    insights: dict = {}
    change_notes: list[str] = []

    # Background style
    if "bg" in changes and changes["bg"]["status"] == "changed":
        new_bg = changes["bg"]["new"]
        if new_bg["x0"] < 0.01 and new_bg["x1"] > 0.99:
            insights["background_style"] = "full_warm"
            change_notes.append("background expanded to full-slide warm fill")
        if "fill" in new_bg:
            insights["background_color"] = new_bg["fill"]

    # Arrow style change (connector → block shape)
    if "arrow" in changes:
        ch = changes["arrow"]
        if ch["status"] == "changed":
            old_type = ch["old"].get("shape_type", "")
            new_type = ch["new"].get("shape_type", "")
            if "CONNECTOR" in old_type and "AUTO_SHAPE" in new_type:
                insights["arrow_style"] = "block"
                rot = ch["new"].get("rotation", 0)
                insights["arrow_block"] = {
                    "x0": ch["new"]["x0"], "y0": ch["new"]["y0"],
                    "x1": ch["new"]["x1"], "y1": ch["new"]["y1"],
                    "rotation": rot,
                }
                change_notes.append(
                    f"arrow replaced: connector → block shape (rotation {rot:.1f}°)"
                )
        elif ch["status"] == "changed":
            # Block arrow repositioned/resized
            insights["arrow_block"] = {
                "x0": ch["new"]["x0"], "y0": ch["new"]["y0"],
                "x1": ch["new"]["x1"], "y1": ch["new"]["y1"],
                "rotation": ch["new"].get("rotation", 0),
            }
            change_notes.append("arrow block repositioned/resized")

    # Chart resize
    if "chart" in changes and changes["chart"]["status"] == "changed":
        new_c = changes["chart"]["new"]
        insights["chart"] = {
            "x0": new_c["x0"], "y0": new_c["y0"],
            "x1": new_c["x1"], "y1": new_c["y1"],
        }
        d = changes["chart"]["delta"]
        change_notes.append(
            f"chart resized: dw={d['dw']:+.3f} dh={d['dh']:+.3f} fraction"
        )

    # Icon position/size corrections
    icon_scale_ratios: list[float] = []
    icon_changes: list[str] = []
    icon_prefs: dict = {}

    for name, ch in changes.items():
        if not name.startswith("icon_"):
            continue
        icon_key = name[len("icon_"):]   # strip "icon_" prefix

        if ch["status"] == "changed":
            new_b = ch["new"]
            icon_prefs[icon_key] = {
                "x0": new_b["x0"], "y0": new_b["y0"],
                "x1": new_b["x1"], "y1": new_b["y1"],
            }
            # Scale ratio: sqrt(area_new / area_old)
            old_b = ch["old"]
            old_area = max(1e-9, (old_b["x1"] - old_b["x0"]) * (old_b["y1"] - old_b["y0"]))
            new_area = max(1e-9, (new_b["x1"] - new_b["x0"]) * (new_b["y1"] - new_b["y0"]))
            scale = math.sqrt(new_area / old_area)
            icon_scale_ratios.append(scale)
            icon_changes.append(f"{icon_key}: scale={scale:.2f}")

    if icon_prefs:
        insights["icons"] = icon_prefs

    if icon_scale_ratios:
        avg = sum(icon_scale_ratios) / len(icon_scale_ratios)
        insights["_avg_icon_scale"] = round(avg, 3)
        if avg < 0.75:
            insights["_icon_size_preference"] = "smaller"
            change_notes.append(
                f"icons consistently shrunk (avg scale {avg:.2f}×) — "
                f"learn smaller defaults: {', '.join(icon_changes)}"
            )
        elif avg > 1.25:
            insights["_icon_size_preference"] = "larger"
            change_notes.append(
                f"icons consistently enlarged (avg scale {avg:.2f}×)"
            )
        else:
            change_notes.append(f"selective icon resizing: {', '.join(icon_changes)}")

    # Text box positions
    text_prefs: dict = {}
    for name, ch in changes.items():
        if name.startswith("text_") and ch["status"] == "changed":
            key = name[len("text_"):]
            new_b = ch["new"]
            text_prefs[key] = {
                "x0": new_b["x0"], "y0": new_b["y0"],
                "x1": new_b["x1"], "y1": new_b["y1"],
            }
            change_notes.append(f"text '{key}' repositioned")
    if text_prefs:
        insights["text"] = text_prefs

    insights["_change_notes"] = change_notes
    return insights


# ── Preferences update ─────────────────────────────────────────────────────────

def _load_prefs() -> dict:
    if PREFS_FILE.exists():
        with open(PREFS_FILE) as f:
            return json.load(f)
    return {"_meta": {"sessions": 0, "last_updated": "", "notes": []}}


def _save_prefs(prefs: dict, insights: dict, session_notes: list[str]) -> None:
    # Bump session counter
    meta = prefs.setdefault("_meta", {"sessions": 0, "last_updated": "", "notes": []})
    meta["sessions"] = meta.get("sessions", 0) + 1
    meta["last_updated"] = str(date.today())

    note = f"Session {meta['sessions']} ({date.today()}): " + "; ".join(session_notes or ["no changes detected"])
    meta.setdefault("notes", []).append(note)

    # Persist style-level insights
    style = prefs.setdefault("style", {})
    if "background_style" in insights:
        style["background"] = insights["background_style"]
    if "background_color" in insights:
        style["background_color"] = insights["background_color"]
    if "arrow_style" in insights:
        style["arrow_style"] = insights["arrow_style"]

    # Persist chart
    if "chart" in insights:
        prefs["chart"] = insights["chart"]

    # Persist arrow block
    if "arrow_block" in insights:
        prefs["arrow_block"] = insights["arrow_block"]

    # Persist icon positions (deep merge: only update icons that changed)
    if "icons" in insights:
        prefs.setdefault("icons", {}).update(insights["icons"])

    # Persist text positions
    if "text" in insights:
        prefs.setdefault("text", {}).update(insights["text"])

    PREFS_FILE.write_text(json.dumps(prefs, indent=2))


# ── Report ─────────────────────────────────────────────────────────────────────

def _print_report(changes: dict, insights: dict) -> None:
    print("\n" + "═" * 60)
    print("  sync_pptx_feedback — diff report")
    print("═" * 60)

    if not changes:
        print("  No changes detected — PPTX matches baseline exactly.")
        return

    # Group by status
    changed = {n: c for n, c in changes.items() if c["status"] == "changed"}
    deleted = {n: c for n, c in changes.items() if c["status"] == "deleted"}
    added   = {n: c for n, c in changes.items() if c["status"] == "added"}

    if changed:
        print(f"\n  Changed shapes ({len(changed)})")
        print("  " + "─" * 56)
        for name, ch in changed.items():
            d = ch.get("delta", {})
            rot_old = ch["old"].get("rotation", 0) or 0
            rot_new = ch["new"].get("rotation", 0) or 0
            rot_str = f"  rot: {rot_old:.1f}°→{rot_new:.1f}°" if abs(rot_new - rot_old) > 0.5 else ""
            type_str = ""
            if ch["old"].get("shape_type") != ch["new"].get("shape_type"):
                type_str = f"  type: {ch['old']['shape_type']}→{ch['new']['shape_type']}"
            print(f"  {name}")
            print(f"    pos:  ({ch['old']['x0']:.3f},{ch['old']['y0']:.3f})→"
                  f"({ch['new']['x0']:.3f},{ch['new']['y0']:.3f}){rot_str}{type_str}")
            if d.get("dw") or d.get("dh"):
                print(f"    size: dw={d['dw']:+.3f}  dh={d['dh']:+.3f}")

    if deleted:
        print(f"\n  Deleted shapes ({len(deleted)}): {', '.join(deleted)}")

    if added:
        print(f"\n  Added shapes ({len(added)}): {', '.join(added)}")

    print("\n  Insights learned")
    print("  " + "─" * 56)
    for note in insights.get("_change_notes", []):
        print(f"  · {note}")
    if "_avg_icon_scale" in insights:
        pref = insights.get("_icon_size_preference", "mixed")
        print(f"  · avg icon scale factor: {insights['_avg_icon_scale']:.2f}×  ({pref})")

    print("\n" + "═" * 60)


# ── Main ───────────────────────────────────────────────────────────────────────

def main() -> None:
    if not PPTX_FILE.exists():
        print(f"Error: PPTX not found at {PPTX_FILE}")
        return
    if not BASELINE_FILE.exists():
        print(f"Error: baseline not found at {BASELINE_FILE}")
        print("Run generate_tariff_scene.py first to create a baseline.")
        return

    print(f"Reading {PPTX_FILE.name}…")
    current  = _extract_shapes(PPTX_FILE)

    print(f"Reading baseline {BASELINE_FILE.name}…")
    with open(BASELINE_FILE) as f:
        baseline = json.load(f)

    changes  = _diff(baseline, current)
    insights = _derive_insights(changes, current)
    prefs    = _load_prefs()

    _print_report(changes, insights)

    notes = insights.pop("_change_notes", [])
    # Remove internal keys before saving
    clean_insights = {k: v for k, v in insights.items() if not k.startswith("_")}
    _save_prefs(prefs, clean_insights, notes)

    print(f"\nPreferences updated → {PREFS_FILE}")
    print("Run generate_tariff_scene.py to regenerate with corrected positions.")


if __name__ == "__main__":
    main()
